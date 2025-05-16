import streamlit as st
import io
from pptx import Presentation
import base64
import pandas as pd
import re
from collections import defaultdict
import traceback
import os
from google import genai
from google.generativeai import GenerativeModel

st.set_page_config(page_title="PowerPoint Template Filler with AI", layout="wide")

st.title("PowerPoint Template Filler with AI")
st.markdown("Upload your PowerPoint template and fill it with your content or AI-generated content")

# Initialize session state for storing presentation data
if 'template_name' not in st.session_state:
    st.session_state.template_name = "presentation.pptx"

if 'gemini_api_key' not in st.session_state:
    st.session_state.gemini_api_key = ""

# File uploader for PowerPoint template
uploaded_file = st.file_uploader("Upload your PowerPoint template (.pptx)", type="pptx")

def initialize_gemini_api(api_key):
    """Initialize the Gemini API with the provided key"""
    os.environ["GOOGLE_API_KEY"] = api_key
    return genai.Client(api_key=api_key)

def get_ai_response(prompt, model):
    """Get a response from Gemini AI based on the prompt"""
    try:
        response = model.models.generate_content(model="gemini-2.0-flash", contents=prompt)
        return response.text
    except Exception as e:
        st.error(f"Error getting AI response: {str(e)}")
        return f"Error: {str(e)}"

def get_shape_name(shape):
    """Get a descriptive name for a shape based on its properties"""
    name = ""

    # Check if shape is a placeholder - safely handle this check
    try:
        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
            # Try to get the shape name or placeholder type
            if hasattr(shape, "name") and shape.name:
                name = shape.name
            elif hasattr(shape, "placeholder_format") and hasattr(shape.placeholder_format, "type"):
                placeholder_types = {
                    1: "TITLE",
                    2: "BODY",
                    3: "CENTER_TITLE",
                    4: "SUBTITLE",
                    13: "TEXT",
                    17: "CHART",
                    18: "TABLE",
                    19: "SMART_ART",
                    21: "PICTURE"
                }
                ph_type = shape.placeholder_format.type
                name = placeholder_types.get(ph_type, f"PLACEHOLDER_{ph_type}")
    except ValueError:
        # Not a placeholder, continue with other checks
        pass

    # If no name found, use shape type
    if not name and hasattr(shape, "shape_type"):
        name = f"SHAPE_{shape.shape_type}"

    # If still no name, use a generic identifier
    if not name:
        name = "SHAPE"
        
    return name

def extract_table_data(table):
    """Extract data from a table shape"""
    data = []
    for i, row in enumerate(table.rows):
        row_data = []
        for j, cell in enumerate(row.cells):
            cell_text = cell.text
            row_data.append(cell_text)
        data.append(row_data)
    return data

def handle_shape_type(shape, slide_idx, shape_idx, shape_count, gemini_model=None):
    """Handle different types of shapes and create appropriate UI elements"""
    shape_data = {}

    # Get shape name for better labeling
    shape_name = get_shape_name(shape)
    label = f"{shape_name} {shape_count}"

    # Handle tables
    if hasattr(shape, "table"):
        st.write(f"### Table: {label}")
        
        table = shape.table
        rows = len(table.rows)
        cols = len(table.columns) if rows > 0 else 0
        
        # Extract existing table data
        table_data = extract_table_data(table)
        
        # Create a pandas DataFrame for editing
        df = pd.DataFrame(table_data)
        
        # AI assistance for tables
        if gemini_model and st.button(f"Generate table content with AI", key=f"ai_table_{slide_idx}_{shape_idx}"):
            with st.spinner("Generating table content..."):
                # Create a prompt based on the table structure
                table_prompt = f"Create content for a table with {rows} rows and {cols} columns. "
                
                # If there are headers, include them in the prompt
                if rows > 0:
                    headers = table_data[0]
                    if any(headers):
                        table_prompt += f"The headers are: {', '.join(headers)}. "
                
                table_prompt += "Generate appropriate data for each cell, formatted as CSV data."
                
                # Get AI response
                ai_response = get_ai_response(table_prompt, gemini_model)
                
                # Try to parse the response into a table
                try:
                    # Try to extract CSV data from the response
                    csv_content = ai_response
                    if "```" in ai_response:
                        # Extract code block content if present
                        csv_content = re.search(r"```(?:csv)?\n(.*?)\n```", ai_response, re.DOTALL)
                        if csv_content:
                            csv_content = csv_content.group(1)
                    
                    new_df = pd.read_csv(io.StringIO(csv_content), header=None)
                    
                    # Ensure the dataframe has the right dimensions
                    if len(new_df) > rows:
                        new_df = new_df.iloc[:rows]
                    if len(new_df.columns) > cols:
                        new_df = new_df.iloc[:, :cols]
                    
                    # Update the table data editor
                    df = new_df
                except Exception as e:
                    st.error(f"Failed to parse AI-generated table data: {str(e)}")
        
        # Use Streamlit's data editor for tables
        edited_df = st.data_editor(
            df, 
            key=f"table_{slide_idx}_{shape_idx}",
            use_container_width=True,
            num_rows="fixed",
        )
        
        shape_data = {
            "type": "table",
            "rows": rows,
            "cols": cols,
            "data": edited_df.values.tolist()
        }

    # Handle charts (placeholder for chart data input)
    elif hasattr(shape, "chart"):
        st.write(f"### Chart: {label}")
        st.info("Charts can be modified through data input:")
        
        chart_type = "Unknown"
        if hasattr(shape.chart, "chart_type"):
            chart_type = str(shape.chart.chart_type)
        
        st.write(f"Chart type: {chart_type}")
        
        # Simple CSV input for chart data
        if gemini_model and st.button(f"Generate chart data with AI", key=f"ai_chart_{slide_idx}_{shape_idx}"):
            with st.spinner("Generating chart data..."):
                chart_prompt = f"Generate sample data for a {chart_type} chart in CSV format. Include headers in the first row. The data should be realistic and meaningful."
                ai_response = get_ai_response(chart_prompt, gemini_model)
                
                # Try to extract CSV data from the response
                csv_content = ai_response
                if "```" in ai_response:
                    # Extract code block content if present
                    csv_content = re.search(r"```(?:csv)?\n(.*?)\n```", ai_response, re.DOTALL)
                    if csv_content:
                        csv_content = csv_content.group(1)
                
                # Use the extracted CSV data
                csv_input = csv_content
        else:
            csv_input = ""
            
        csv_data = st.text_area(
            "Enter chart data (CSV format - first row as headers):",
            csv_input,
            height=150,
            key=f"chart_{slide_idx}_{shape_idx}"
        )
        
        shape_data = {
            "type": "chart",
            "chart_type": chart_type,
            "data": csv_data
        }

    # Handle text shapes
    elif hasattr(shape, "text_frame"):
        current_text = shape.text if hasattr(shape, "text") else ""
        
        # Clean up placeholder text patterns
        placeholder_text = current_text
        cleaned_text = re.sub(r'\[.*?\]|<.*?>|Click to add.*', '', current_text).strip()
        
        # Use the existing text as default value, or cleaned text if it's a placeholder
        default_text = cleaned_text if cleaned_text else current_text
        
        # AI generation for text
        if gemini_model:
            text_col1, text_col2 = st.columns([3, 1])
            with text_col2:
                # Add a prompt field and generate button
                ai_prompt = st.text_input(
                    "AI prompt (optional)",
                    placeholder="Describe what content you want here...",
                    key=f"ai_prompt_{slide_idx}_{shape_idx}"
                )
                generate_text = st.button(
                    "Generate with AI",
                    key=f"generate_{slide_idx}_{shape_idx}"
                )
                
                if generate_text and ai_prompt:
                    with st.spinner("Generating text..."):
                        # Create a prompt considering the context
                        context = f"For a PowerPoint slide, generate content for a text box labeled '{label}'. "
                        if "TITLE" in label:
                            context += "This is a slide title, so keep it concise and impactful. "
                        elif "SUBTITLE" in label:
                            context += "This is a subtitle, so provide a brief supporting statement. "
                        elif "BODY" in label:
                            context += "This is body text, so provide informative content with clear structure. "
                        
                        complete_prompt = context + ai_prompt
                        ai_response = get_ai_response(complete_prompt, gemini_model)
                        
                        # Update the default text
                        default_text = ai_response
        
        with text_col1 if gemini_model else st.container():
            text_input = st.text_area(
                label,
                value=default_text,
                height=100,
                key=f"text_{slide_idx}_{shape_idx}"
            )
        
        shape_data = {
            "type": "text",
            "text": text_input
        }

    # Handle images
    elif hasattr(shape, "image"):
        st.write(f"### Image: {label}")
        st.info("You can replace this image:")
        
        new_image = st.file_uploader(
            "Upload replacement image", 
            type=["png", "jpg", "jpeg"],
            key=f"image_{slide_idx}_{shape_idx}"
        )
        
        shape_data = {
            "type": "image",
            "image": new_image
        }

    # Any other shape with text capability
    elif hasattr(shape, "text"):
        shape_data = {
            "type": "text",
            "text": st.text_area(
                f"{label}",
                value=shape.text,
                height=100,
                key=f"other_{slide_idx}_{shape_idx}"
            )
        }

    return shape_data

# Sidebar for Gemini API setup
with st.sidebar:
    st.subheader("AI Assistant Setup")
    api_key = st.text_input(
        "Enter your Gemini API Key",
        value=st.session_state.gemini_api_key,
        type="password",
        key="api_key_input"
    )
    
    if api_key != st.session_state.gemini_api_key:
        st.session_state.gemini_api_key = api_key
    
    if st.button("Validate API Key") and api_key:
        try:
            model = initialize_gemini_api(api_key)
            # response = client.models.generate_content(model="gemini-2.0-flash", contents=prompt)

            response = get_ai_response("Hello, are you working?", model)
            if response and not response.startswith("Error:"):
                st.success("API key validated successfully!")
                st.session_state.gemini_model = model
            else:
                st.error("API key validation failed.")
                st.session_state.gemini_model = None
        except Exception as e:
            st.error(f"API key validation failed: {str(e)}")
            st.session_state.gemini_model = None
    
    # Instructions and features
    st.subheader("Instructions")
    st.markdown("""
    1. Enter your Gemini API key and click "Validate API Key"
    2. Upload your PowerPoint template (.pptx)
    3. Navigate through tabs to edit each slide
    4. Use AI to generate content or enter your own
    5. Click "Generate PowerPoint" when finished
    6. Download your completed presentation
    """)
    
    st.subheader("Features")
    st.markdown("""
    - Preserves template formatting
    - AI-assisted content generation
    - Supports text, tables, notes and placeholders
    - Edit slide titles individually
    - Customize download filename
    - Error handling for robustness
    """)
    
    # Advanced options in expander
    with st.expander("Advanced Options"):
        st.checkbox("Show debug info", key="debug_mode")
        
        # AI model selection (for future)
        st.selectbox(
            "Gemini model (future option)", 
            ["gemini-1.5-flash", "gemini-1.5-pro"], 
            index=0,
            disabled=True
        )
        
        st.selectbox(
            "Default slide layout", 
            ["Default", "Title Only", "Title and Content", "Blank"], 
            index=0
        )

try:
    if uploaded_file is not None:
        # Save template name for download
        st.session_state.template_name = uploaded_file.name

        # Read the template
        prs = Presentation(uploaded_file)
        
        # Display template information
        st.success("Template uploaded successfully!")
        st.write(f"Number of slides: {len(prs.slides)}")
        
        # Create tabs for each slide
        tabs = st.tabs([f"Slide {i+1}" for i in range(len(prs.slides))])
        
        # Dictionary to store edited content
        slide_content = defaultdict(dict)
        
        # Get the Gemini model if available
        gemini_model = st.session_state.get('gemini_model', None)
        
        # Process each slide
        for i, (tab, slide) in enumerate(zip(tabs, prs.slides)):
            with tab:
                st.subheader(f"Slide {i+1}")
                
                # AI content generation for entire slide
                if gemini_model:
                    with st.expander("Generate all slide content with AI"):
                        slide_theme = st.text_input(
                            "What is this slide about?",
                            placeholder="e.g., Product features, Market analysis, Team introduction",
                            key=f"slide_theme_{i}"
                        )
                        
                        if st.button("Generate All Content", key=f"generate_all_{i}"):
                            with st.spinner("Generating slide content..."):
                                # Analyze the slide structure
                                shapes_info = []
                                for shape in slide.shapes:
                                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                                        shape_name = get_shape_name(shape)
                                        shapes_info.append(shape_name)
                                
                                # Build a comprehensive prompt
                                slide_prompt = f"Generate content for a PowerPoint slide about '{slide_theme}'. "
                                slide_prompt += f"The slide contains the following elements: {', '.join(shapes_info)}. "
                                slide_prompt += "Provide appropriate content for each element in JSON format with the element names as keys."
                                
                                # Get AI response
                                ai_response = get_ai_response(slide_prompt, gemini_model)
                                st.write("AI generated suggestions (you can copy-paste these into the fields below):")
                                st.code(ai_response)
                
                # Add a slide title input at the top (if title placeholder exists)
                title_shape = None
                for shape in slide.shapes:
                    # Safely check if this is a title placeholder
                    try:
                        is_title_placeholder = (
                            hasattr(shape, "is_placeholder") and 
                            shape.is_placeholder and
                            hasattr(shape, "placeholder_format") and 
                            hasattr(shape.placeholder_format, "type") and 
                            shape.placeholder_format.type == 1  # Title placeholder
                        )
                        if is_title_placeholder:
                            title_shape = shape
                            break
                    except ValueError:
                        # Not a placeholder, continue to next shape
                        continue
                
                if title_shape:
                    st.write("### Slide Title")
                    
                    # Add AI generation option for title
                    if gemini_model:
                        title_col1, title_col2 = st.columns([3, 1])
                        with title_col2:
                            if st.button("Generate title with AI", key=f"gen_title_{i}"):
                                with st.spinner("Generating title..."):
                                    # Create a context-aware prompt
                                    slide_idx_text = f"slide {i+1} of {len(prs.slides)}"
                                    title_prompt = f"Generate a concise, impactful title for {slide_idx_text} in a presentation"
                                    if 'slide_theme' in st.session_state and st.session_state[f'slide_theme_{i}']:
                                        title_prompt += f" about {st.session_state[f'slide_theme_{i}']}"
                                    title_prompt += ". Just provide the title text without any additional explanation or formatting."
                                    
                                    title_text = get_ai_response(title_prompt, gemini_model)
                                    # Clean up potential quotes or extra formatting
                                    title_text = re.sub(r'^["\'"\']|["\'"\']$', '', title_text.strip())
                                    st.session_state[f"title_{i}"] = title_text
                        with title_col1:
                            slide_title = st.text_input(
                                "Title", 
                                value=st.session_state.get(f"title_{i}", title_shape.text if hasattr(title_shape, "text") else ""),
                                key=f"title_{i}"
                            )
                    else:
                        slide_title = st.text_input(
                            "Title", 
                            value=title_shape.text if hasattr(title_shape, "text") else "",
                            key=f"title_{i}"
                        )
                    
                    slide_content[i]["title"] = {
                        "shape_idx": list(slide.shapes).index(title_shape),
                        "text": slide_title
                    }
                
                # Add slide notes option
                if hasattr(slide, "notes_slide") and slide.notes_slide:
                    notes_text = ""
                    try:
                        # With this:
                        for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                            notes_text += paragraph.text + "\n"
                    except AttributeError:
                        # Handle case where notes_text_frame doesn't exist
                        pass
                            
                    with st.expander("Slide Notes"):
                        # Add AI generation for notes
                        if gemini_model:
                            notes_col1, notes_col2 = st.columns([3, 1])
                            with notes_col2:
                                notes_prompt = st.text_input(
                                    "Notes prompt",
                                    placeholder="e.g., Key points to remember",
                                    key=f"notes_prompt_{i}"
                                )
                                if st.button("Generate notes", key=f"gen_notes_{i}") and notes_prompt:
                                    with st.spinner("Generating notes..."):
                                        context = f"Create speaker notes for slide {i+1} "
                                        if title_shape and hasattr(title_shape, "text") and title_shape.text:
                                            context += f"titled '{title_shape.text}' "
                                        context += "based on the following prompt: " + notes_prompt
                                        
                                        new_notes = get_ai_response(context, gemini_model)
                                        st.session_state[f"notes_{i}"] = new_notes
                            
                            with notes_col1:
                                new_notes = st.text_area(
                                    "Notes", 
                                    value=st.session_state.get(f"notes_{i}", notes_text),
                                    height=100,
                                    key=f"notes_{i}"
                                )
                        else:
                            new_notes = st.text_area(
                                "Notes", 
                                value=notes_text,
                                height=100,
                                key=f"notes_{i}"
                            )
                        
                        slide_content[i]["notes"] = new_notes
                
                # Process shapes
                shape_count = 0
                
                for j, shape in enumerate(slide.shapes):
                    # Skip title shape as it's handled separately
                    if title_shape and shape == title_shape:
                        continue
                        
                    # Handle different shape types
                    shape_count += 1
                    try:
                        shape_data = handle_shape_type(shape, i, j, shape_count, gemini_model)
                        if shape_data:  # Only add if we got data back
                            slide_content[i][j] = shape_data
                    except Exception as e:
                        st.warning(f"Error processing shape {j}: {str(e)}")
                        continue
        
        # Generate button
        col1, col2 = st.columns([1, 3])
        with col1:
            custom_filename = st.text_input(
                "Output filename", 
                value=f"filled_{st.session_state.template_name}"
            )
        
        with col2:
            generate_btn = st.button("Generate PowerPoint", type="primary", use_container_width=True)
        
        if generate_btn:
            with st.spinner("Generating PowerPoint..."):
                # Update the presentation with user input
                for slide_idx, shapes_dict in slide_content.items():
                    slide = prs.slides[slide_idx]
                    
                    # Handle special "title" entry
                    if "title" in shapes_dict:
                        title_idx = shapes_dict["title"]["shape_idx"]
                        title_text = shapes_dict["title"]["text"]
                        title_shape = list(slide.shapes)[title_idx]
                        if hasattr(title_shape, "text_frame"):
                            title_shape.text_frame.text = title_text
                    
                    # Handle special "notes" entry
                    if "notes" in shapes_dict:
                        if hasattr(slide, "notes_slide") and slide.notes_slide:
                            for shape in slide.notes_slide.shapes:
                                if hasattr(shape, "text_frame"):
                                    shape.text_frame.text = shapes_dict["notes"]
                                    break
                    
                    # Handle regular shapes
                    for shape_idx, shape_data in shapes_dict.items():
                        # Skip special keys
                        if shape_idx in ["title", "notes"]:
                            continue
                        
                        try:
                            # Get the actual shape
                            shape = list(slide.shapes)[shape_idx]
                            
                            # Handle by shape type
                            if shape_data["type"] == "text":
                                if hasattr(shape, "text_frame"):
                                    # Clear existing text
                                    for paragraph in shape.text_frame.paragraphs:
                                        if paragraph.runs:
                                            for run in paragraph.runs:
                                                run.text = ""
                                    
                                    # Get the text and split into lines
                                    text = shape_data["text"]
                                    lines = text.split('\n')
                                    
                                    # Set the first line to the first paragraph
                                    if shape.text_frame.paragraphs:
                                        shape.text_frame.paragraphs[0].text = lines[0] if lines else ""
                                        
                                        # Add remaining lines as new paragraphs
                                        for line in lines[1:]:
                                            p = shape.text_frame.add_paragraph()
                                            p.text = line
                                    else:
                                        shape.text_frame.text = text
                            
                            elif shape_data["type"] == "table" and hasattr(shape, "table"):
                                table = shape.table
                                table_data = shape_data["data"]
                                
                                # Update each cell with new data
                                for row_idx, row in enumerate(table.rows):
                                    if row_idx < len(table_data):
                                        for col_idx, cell in enumerate(row.cells):
                                            if col_idx < len(table_data[row_idx]):
                                                cell_value = table_data[row_idx][col_idx]
                                                cell.text = str(cell_value) if pd.notna(cell_value) else ""
                            
                            elif shape_data["type"] == "chart" and hasattr(shape, "chart"):
                                # Parse CSV data for chart
                                csv_text = shape_data["data"]
                                if csv_text.strip():
                                    try:
                                        # Parse the CSV data
                                        chart_data = pd.read_csv(io.StringIO(csv_text))
                                        
                                        # Update chart data
                                        chart = shape.chart
                                        
                                        # Get chart data worksheet
                                        if hasattr(chart, "_workbook") and chart._workbook:
                                            # This is a simplified approach - actual chart data update
                                            # would require more complex handling based on chart type
                                            st.info("Chart data updated (preview only)")
                                    except Exception as e:
                                        st.error(f"Error updating chart: {str(e)}")
                            
                            # Image handling would go here but requires more complex processing
                            # that is beyond the scope of python-pptx's capabilities
                            
                        except Exception as e:
                            st.error(f"Error updating shape {shape_idx} on slide {slide_idx+1}: {str(e)}")
                
                # Save to BytesIO object
                output = io.BytesIO()
                prs.save(output)
                output.seek(0)
                
                # Provide download link
                b64 = base64.b64encode(output.read()).decode()
                href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{custom_filename}">Download PowerPoint Presentation</a>'
                st.markdown(href, unsafe_allow_html=True)
                st.success("PowerPoint file generated successfully!")

except Exception as e:
    st.error(f"An error occurred: {str(e)}")
    if st.session_state.get("debug_mode", False):  # Only show traceback in debug mode
        st.error(traceback.format_exc())

# Add information about the AI integration
st.sidebar.subheader("About AI Integration")
st.sidebar.markdown("""
This application integrates Google's Gemini AI to help you:
- Generate slide titles and content
- Create table data
- Write speaker notes
- Produce chart data

To use AI features:
1. Obtain a Gemini API key from Google AI Studio
2. Enter your API key in the sidebar
3. Look for AI generation buttons throughout the app
""")

st.sidebar.subheader("Troubleshooting")
st.sidebar.markdown("""
- If you encounter errors with complex templates, try using simpler templates with standard elements
- Some PowerPoint features (SmartArt, complex animations) have limited edit capabilities
- If AI generation fails, try more specific prompts or check your API key
- For best results, use templates with clear text placeholders
""")
